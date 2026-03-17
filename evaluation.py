from __future__ import annotations

from collections import Counter
from dataclasses import dataclass
from pathlib import Path
from typing import Any
import math
import xml.etree.ElementTree as ET

from pptx import Presentation

from .models import BBox, FigureSceneGraph, HierarchyUnit
from .utils import load_json, normalize_text, polyline_length, save_json, similarity


EDITABLE_KINDS = {"text", "container", "icon", "group", "background"}


@dataclass
class EvaluationResult:
    case_id: str
    method_name: str
    metrics: dict[str, float]
    failure_tags: list[str]
    metadata: dict[str, Any]

    def to_dict(self) -> dict[str, Any]:
        return {
            "case_id": self.case_id,
            "method_name": self.method_name,
            "metrics": self.metrics,
            "failure_tags": self.failure_tags,
            "metadata": self.metadata,
        }


def _safe_load_scene(path: str | Path | None) -> FigureSceneGraph | None:
    if not path:
        return None
    path = Path(path)
    if not path.exists():
        return None
    return FigureSceneGraph.load(path)


def _safe_load_hierarchy(path: str | Path | None) -> dict[str, Any] | None:
    if not path:
        return None
    path = Path(path)
    if not path.exists():
        return None
    return load_json(path)


def _label_counter(scene: FigureSceneGraph) -> Counter[str]:
    return Counter(normalize_text(node.label) for node in scene.nodes if node.kind != "background")


def _multiset_f1(pred: Counter[str], target: Counter[str]) -> tuple[float, float, float]:
    if not pred and not target:
        return 1.0, 1.0, 1.0
    correct = sum((pred & target).values())
    precision = correct / max(1, sum(pred.values()))
    recall = correct / max(1, sum(target.values()))
    f1 = 0.0 if precision + recall == 0 else 2 * precision * recall / (precision + recall)
    return precision, recall, f1


def _edge_counter(scene: FigureSceneGraph) -> Counter[str]:
    node_map = scene.node_map()
    edge_labels: Counter[str] = Counter()
    for edge in scene.edges:
        source = node_map.get(edge.source_id)
        target = node_map.get(edge.target_id)
        if not source or not target:
            continue
        edge_labels[f"{normalize_text(source.label)}->{normalize_text(target.label)}"] += 1
    return edge_labels


def _group_score(pred: FigureSceneGraph, target: FigureSceneGraph) -> float:
    if not target.groups:
        return 1.0
    pred_groups = {normalize_text(group.label): group for group in pred.groups}
    pred_nodes = pred.node_map()
    scores: list[float] = []
    for target_group in target.groups:
        key = normalize_text(target_group.label)
        matched = pred_groups.get(key)
        if not matched:
            scores.append(0.0)
            continue
        target_members = {normalize_text(target.node_map()[node_id].label) for node_id in target_group.node_ids if node_id in target.node_map()}
        pred_members = {normalize_text(pred_nodes[node_id].label) for node_id in matched.node_ids if node_id in pred_nodes}
        union = target_members | pred_members
        scores.append(1.0 if not union else len(target_members & pred_members) / len(union))
    return sum(scores) / len(scores)


def _levenshtein(a: list[str], b: list[str]) -> int:
    if len(a) < len(b):
        a, b = b, a
    previous = list(range(len(b) + 1))
    for i, token_a in enumerate(a, start=1):
        current = [i]
        for j, token_b in enumerate(b, start=1):
            insertions = previous[j] + 1
            deletions = current[j - 1] + 1
            substitutions = previous[j - 1] + (0 if token_a == token_b else 1)
            current.append(min(insertions, deletions, substitutions))
        previous = current
    return previous[-1]


def _greedy_label_matches(pred: FigureSceneGraph, target: FigureSceneGraph) -> list[tuple[str, str, float]]:
    pred_nodes = [node for node in pred.nodes if node.kind != "background"]
    target_nodes = [node for node in target.nodes if node.kind != "background"]
    matches: list[tuple[str, str, float]] = []
    used_pred: set[str] = set()
    for target_node in target_nodes:
        best_node = None
        best_score = -1.0
        for pred_node in pred_nodes:
            if pred_node.id in used_pred:
                continue
            score = similarity(pred_node.label, target_node.label)
            if score > best_score:
                best_score = score
                best_node = pred_node
        if best_node is not None:
            used_pred.add(best_node.id)
            matches.append((best_node.id, target_node.id, best_score))
    return matches


def _text_metrics(pred: FigureSceneGraph, target: FigureSceneGraph) -> dict[str, float]:
    matches = _greedy_label_matches(pred, target)
    pred_map = pred.node_map()
    target_map = target.node_map()
    exact_hits = 0
    total = 0
    pred_chars: list[str] = []
    target_chars: list[str] = []
    pred_words: list[str] = []
    target_words: list[str] = []
    for pred_id, target_id, _ in matches:
        pred_label = normalize_text(pred_map[pred_id].label)
        target_label = normalize_text(target_map[target_id].label)
        total += 1
        if pred_label == target_label:
            exact_hits += 1
        pred_chars.extend(list(pred_label))
        target_chars.extend(list(target_label))
        pred_words.extend(pred_label.split())
        target_words.extend(target_label.split())
    exact_match = exact_hits / max(1, total)
    cer = _levenshtein(pred_chars, target_chars) / max(1, len(target_chars))
    wer = _levenshtein(pred_words, target_words) / max(1, len(target_words))
    return {
        "text_exact_match": round(exact_match, 4),
        "ocr_cer": round(cer, 4),
        "ocr_wer": round(wer, 4),
        "text_accuracy_score": round((exact_match + (1 - cer) + (1 - wer)) / 3.0, 4),
    }


def _bbox_similarity(a: BBox, b: BBox, canvas_w: int, canvas_h: int) -> float:
    diagonal = math.dist((0, 0), (canvas_w, canvas_h))
    center_dist = math.dist((a.center_x, a.center_y), (b.center_x, b.center_y)) / max(1.0, diagonal)
    size_penalty = abs(a.width - b.width) / max(1.0, b.width) + abs(a.height - b.height) / max(1.0, b.height)
    return max(0.0, 1.0 - min(1.0, center_dist + 0.5 * size_penalty))


def _pairwise_gap(a: BBox, b: BBox) -> float:
    dx = max(0.0, max(a.x - b.right, b.x - a.right))
    dy = max(0.0, max(a.y - b.bottom, b.y - a.bottom))
    return max(dx, dy)


def _layout_metrics(pred: FigureSceneGraph, target: FigureSceneGraph) -> dict[str, float]:
    pred_nodes = [node for node in pred.nodes if node.kind != "background"]
    target_map = target.node_map()
    pred_map = pred.node_map()
    pairs = 0
    overlaps = 0
    min_gap = 9999.0
    for idx, node in enumerate(pred_nodes):
        for other in pred_nodes[idx + 1 :]:
            pairs += 1
            if node.bbox.overlaps(other.bbox):
                overlaps += 1
            min_gap = min(min_gap, _pairwise_gap(node.bbox, other.bbox))
    overlap_rate = 0.0 if pairs == 0 else overlaps / pairs
    overlap_score = 1.0 - overlap_rate

    bbox_scores: list[float] = []
    for pred_id, target_id, _ in _greedy_label_matches(pred, target):
        bbox_scores.append(_bbox_similarity(pred_map[pred_id].bbox, target_map[target_id].bbox, target.width, target.height))
    alignment_score = sum(bbox_scores) / len(bbox_scores) if bbox_scores else 0.0

    routing_checks = 0
    routing_hits = 0
    for edge in pred.edges:
        if len(edge.points) >= 2 and edge.source_port and edge.target_port:
            routing_checks += 1
            if polyline_length(edge.points) > 0:
                routing_hits += 1
    routing_score = 1.0 if routing_checks == 0 else routing_hits / routing_checks

    min_spacing_score = min(1.0, max(0.0, min_gap / 24.0)) if min_gap != 9999.0 else 1.0
    readability = (overlap_score + alignment_score + routing_score + min_spacing_score) / 4.0
    return {
        "overlap_rate": round(overlap_rate, 4),
        "alignment_score": round(alignment_score, 4),
        "routing_score": round(routing_score, 4),
        "min_spacing_score": round(min_spacing_score, 4),
        "layout_readability_score": round(readability, 4),
    }


def _svg_roundtrip_score(path: Path, labels: list[str]) -> float:
    if not path.exists():
        return 0.0
    try:
        tree = ET.parse(path)
        text_blob = "".join(element.text or "" for element in tree.getroot().iter())
    except ET.ParseError:
        return 0.0
    hits = sum(1 for label in labels if normalize_text(label) in normalize_text(text_blob))
    return hits / max(1, len(labels))


def _pptx_roundtrip_score(path: Path, labels: list[str]) -> float:
    if not path.exists():
        return 0.0
    try:
        presentation = Presentation(path)
    except Exception:
        return 0.0
    text_blob = []
    shape_count = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            shape_count += 1
            if hasattr(shape, "text") and shape.text:
                text_blob.append(shape.text)
    if shape_count == 0:
        return 0.0
    joined = normalize_text(" ".join(text_blob))
    label_hits = sum(1 for label in labels if normalize_text(label) in joined)
    return label_hits / max(1, len(labels))


def _editability_metrics(pred: FigureSceneGraph, target: FigureSceneGraph, result_dir: str | Path) -> dict[str, float]:
    result_dir = Path(result_dir)
    svg_path = result_dir / "editable.svg"
    pptx_path = result_dir / "editable.pptx"
    export_success = float(svg_path.exists() or pptx_path.exists())
    target_editable = len([node for node in target.nodes if node.kind in EDITABLE_KINDS]) + len(target.groups)
    pred_editable = len([node for node in pred.nodes if node.kind in EDITABLE_KINDS]) + len(pred.groups)
    editable_recall = pred_editable / max(1, target_editable)
    labels = [node.label for node in target.nodes if node.kind != "background"]
    roundtrip_svg = _svg_roundtrip_score(svg_path, labels)
    roundtrip_pptx = _pptx_roundtrip_score(pptx_path, labels)
    roundtrip = max(roundtrip_svg, roundtrip_pptx) if export_success else 0.0
    editability = (export_success + min(1.0, editable_recall) + roundtrip) / 3.0
    return {
        "export_success": round(export_success, 4),
        "editable_element_recall": round(min(1.0, editable_recall), 4),
        "roundtrip_edit_survival": round(roundtrip, 4),
        "editability_score": round(editability, 4),
    }


def _consistency_metrics(pred: FigureSceneGraph) -> dict[str, float]:
    node_count = max(1, len([node for node in pred.nodes if node.kind != "background"]))
    semantic_coverage = sum(1 for node in pred.nodes if node.semantic_anchor_ids) / node_count
    layout_coverage = sum(1 for node in pred.nodes if node.layout_anchor_ids) / node_count
    score = (semantic_coverage + layout_coverage) / 2.0
    return {
        "semantic_anchor_coverage": round(semantic_coverage, 4),
        "layout_anchor_coverage": round(layout_coverage, 4),
        "consistency_score": round(score, 4),
    }


def _node_label_lookup(scene: FigureSceneGraph) -> dict[str, str]:
    return {
        node.id: normalize_text(node.label)
        for node in scene.nodes
        if node.kind != "background"
    }


def _hierarchy_units_from_payload(payload: dict[str, Any] | None, level: str) -> list[HierarchyUnit]:
    if not payload:
        return []
    return [
        HierarchyUnit.from_dict(unit)
        for unit in payload.get("hierarchy_units", [])
        if unit.get("level") == level
    ]


def _unit_node_label_sets(scene: FigureSceneGraph, units: list[HierarchyUnit]) -> list[set[str]]:
    node_lookup = _node_label_lookup(scene)
    label_sets: list[set[str]] = []
    for unit in units:
        labels = {node_lookup[node_id] for node_id in unit.node_ids if node_id in node_lookup}
        if labels:
            label_sets.append(labels)
    return label_sets


def _best_set_iou(target_set: set[str], pred_sets: list[set[str]]) -> float:
    if not pred_sets:
        return 0.0
    best = 0.0
    for pred_set in pred_sets:
        union = target_set | pred_set
        if not union:
            best = max(best, 1.0)
        else:
            best = max(best, len(target_set & pred_set) / len(union))
    return best


def _hierarchy_metrics(
    pred: FigureSceneGraph,
    target: FigureSceneGraph,
    target_hierarchy: dict[str, Any] | None,
) -> dict[str, float]:
    pred_module_sets = _unit_node_label_sets(pred, pred.hierarchy_units_by_level("module"))
    pred_region_sets = _unit_node_label_sets(pred, pred.hierarchy_units_by_level("region"))
    pred_block_sets = _unit_node_label_sets(pred, pred.hierarchy_units_by_level("block"))

    if target_hierarchy:
        target_module_sets = _unit_node_label_sets(target, _hierarchy_units_from_payload(target_hierarchy, "module"))
        target_region_sets = _unit_node_label_sets(target, _hierarchy_units_from_payload(target_hierarchy, "region"))
        target_block_sets = _unit_node_label_sets(target, _hierarchy_units_from_payload(target_hierarchy, "block"))
    else:
        target_module_sets = _unit_node_label_sets(target, target.hierarchy_units_by_level("module"))
        target_region_sets = _unit_node_label_sets(target, target.hierarchy_units_by_level("region"))
        target_block_sets = _unit_node_label_sets(target, target.hierarchy_units_by_level("block"))

    def average_best(target_sets: list[set[str]], pred_sets: list[set[str]]) -> float:
        if not target_sets:
            return 1.0
        return sum(_best_set_iou(target_set, pred_sets) for target_set in target_sets) / len(target_sets)

    module_acc = average_best(target_module_sets, pred_module_sets)
    region_acc = average_best(target_region_sets, pred_region_sets)
    block_acc = average_best(target_block_sets, pred_block_sets)
    hierarchy_consistency = (module_acc + region_acc + block_acc) / 3.0
    return {
        "module_assignment_accuracy": round(module_acc, 4),
        "region_boundary_accuracy": round(region_acc, 4),
        "block_grouping_accuracy": round(block_acc, 4),
        "hierarchy_consistency_score": round(hierarchy_consistency, 4),
    }


def classify_failures(metrics: dict[str, float]) -> list[str]:
    tags: list[str] = []
    if metrics.get("text_accuracy_score", 1.0) < 0.8:
        tags.append("文本恢复失败")
    if metrics.get("edge_f1", 1.0) < 0.85:
        tags.append("边连接错误")
    if metrics.get("group_containment_accuracy", 1.0) < 0.85:
        tags.append("组边界错误")
    if metrics.get("alignment_score", 1.0) < 0.75:
        tags.append("跨panel或布局对齐失败")
    if metrics.get("export_success", 1.0) < 1.0 or metrics.get("roundtrip_edit_survival", 1.0) < 0.75:
        tags.append("可编辑导出失败")
    if metrics.get("layout_readability_score", 1.0) < 0.78 and metrics.get("consistency_score", 1.0) < 0.5:
        tags.append("风格与结构不稳定")
    if metrics.get("module_assignment_accuracy", 1.0) < 0.75:
        tags.append("模块层组织失败")
    if metrics.get("region_boundary_accuracy", 1.0) < 0.75:
        tags.append("分区层边界失败")
    if metrics.get("block_grouping_accuracy", 1.0) < 0.75:
        tags.append("区块层关系失败")
    if metrics.get("hierarchy_consistency_score", 1.0) < 0.72:
        tags.append("全图层规划失败")
    return tags


def evaluate_case(
    case_id: str,
    method_name: str,
    result_dir: str | Path,
    target_scene_path: str | Path,
    target_hierarchy_path: str | Path | None = None,
) -> EvaluationResult:
    result_dir = Path(result_dir)
    pred_scene = _safe_load_scene(result_dir / "scene_graph.json")
    if pred_scene is None:
        raise FileNotFoundError(f"Missing predicted scene graph under {result_dir}")
    target_scene = _safe_load_scene(target_scene_path)
    if target_scene is None:
        raise FileNotFoundError(f"Missing target scene graph: {target_scene_path}")
    target_hierarchy = _safe_load_hierarchy(target_hierarchy_path)

    node_precision, node_recall, node_f1 = _multiset_f1(_label_counter(pred_scene), _label_counter(target_scene))
    edge_precision, edge_recall, edge_f1 = _multiset_f1(_edge_counter(pred_scene), _edge_counter(target_scene))
    group_accuracy = _group_score(pred_scene, target_scene)
    text_metrics = _text_metrics(pred_scene, target_scene)
    layout_metrics = _layout_metrics(pred_scene, target_scene)
    editability_metrics = _editability_metrics(pred_scene, target_scene, result_dir)
    consistency_metrics = _consistency_metrics(pred_scene)
    hierarchy_metrics = _hierarchy_metrics(pred_scene, target_scene, target_hierarchy)

    structure_fidelity = node_f1 * 0.4 + edge_f1 * 0.4 + group_accuracy * 0.2
    overall_auto_score = (
        structure_fidelity * 0.30
        + text_metrics["text_accuracy_score"] * 0.20
        + layout_metrics["layout_readability_score"] * 0.20
        + editability_metrics["editability_score"] * 0.20
        + ((consistency_metrics["consistency_score"] + hierarchy_metrics["hierarchy_consistency_score"]) / 2.0) * 0.10
    )
    publication_readiness = (
        text_metrics["text_accuracy_score"] * 0.25
        + layout_metrics["layout_readability_score"] * 0.35
        + editability_metrics["editability_score"] * 0.25
        + structure_fidelity * 0.15
    )

    metrics = {
        "node_precision": round(node_precision, 4),
        "node_recall": round(node_recall, 4),
        "node_f1": round(node_f1, 4),
        "edge_precision": round(edge_precision, 4),
        "edge_recall": round(edge_recall, 4),
        "edge_f1": round(edge_f1, 4),
        "group_containment_accuracy": round(group_accuracy, 4),
        "structure_fidelity": round(structure_fidelity, 4),
        **text_metrics,
        **layout_metrics,
        **editability_metrics,
        **consistency_metrics,
        **hierarchy_metrics,
        "publication_readiness": round(publication_readiness, 4),
        "overall_auto_score": round(overall_auto_score, 4),
    }
    failures = classify_failures(metrics)
    return EvaluationResult(
        case_id=case_id,
        method_name=method_name,
        metrics=metrics,
        failure_tags=failures,
        metadata={"result_dir": str(result_dir), "target_scene_path": str(target_scene_path)},
    )


def save_evaluation(result: EvaluationResult, output_path: str | Path) -> Path:
    return save_json(output_path, result.to_dict())
