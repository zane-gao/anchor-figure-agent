from __future__ import annotations

from pathlib import Path
from xml.sax.saxutils import escape
import math

from PIL import Image, ImageColor, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.util import Emu

from .models import BBox, FigureSceneGraph, SceneEdge, SceneGroup, SceneNode


EMU_PER_PIXEL = 9525


def _load_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "arial.ttf",
        "Arial.ttf",
        "DejaVuSans.ttf",
        "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/msyh.ttc",
    ]
    for candidate in candidates:
        try:
            return ImageFont.truetype(candidate, size=size)
        except OSError:
            continue
    return ImageFont.load_default()


def _hex_to_rgb(value: str) -> tuple[int, int, int]:
    return ImageColor.getrgb(value)


def _hex_to_rgb_color(value: str) -> RGBColor:
    r, g, b = _hex_to_rgb(value)
    return RGBColor(r, g, b)


def _wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont, max_width: int) -> list[str]:
    words = text.split()
    if not words:
        return [text]
    lines: list[str] = []
    current = words[0]
    for word in words[1:]:
        tentative = f"{current} {word}"
        width = draw.textbbox((0, 0), tentative, font=font)[2]
        if width <= max_width:
            current = tentative
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return lines


def _node_style(scene: FigureSceneGraph, node: SceneNode) -> dict[str, str]:
    palette = scene.style_tokens.palette
    style = {
        "fill": palette["node_fill"],
        "stroke": palette["node_stroke"],
        "text": palette["text"],
        "accent": palette["accent"],
    }
    if node.kind == "background":
        style["fill"] = palette["background"]
        style["stroke"] = palette["background"]
    if node.kind == "icon":
        style["fill"] = palette["group_fill"]
    if node.kind == "text":
        style["fill"] = palette["panel"]
    style.update(node.style)
    return style


def _group_style(scene: FigureSceneGraph, group: SceneGroup) -> dict[str, str]:
    palette = scene.style_tokens.palette
    style = {
        "fill": palette["group_fill"],
        "stroke": palette["group_stroke"],
        "text": palette["muted_text"],
    }
    style.update(group.style)
    return style


def orthogonal_route(source: BBox, target: BBox, direction: str = "right") -> list[tuple[float, float]]:
    if direction == "down":
        start = (source.center_x, source.bottom)
        end = (target.center_x, target.y)
        mid_y = (start[1] + end[1]) / 2.0
        return [start, (start[0], mid_y), (end[0], mid_y), end]
    start = (source.right, source.center_y)
    end = (target.x, target.center_y)
    if target.x < source.x:
        start = (source.center_x, source.bottom)
        end = (target.center_x, target.y)
        mid_y = max(source.bottom + 30, (start[1] + end[1]) / 2.0)
        return [start, (start[0], mid_y), (end[0], mid_y), end]
    mid_x = (start[0] + end[0]) / 2.0
    return [start, (mid_x, start[1]), (mid_x, end[1]), end]


def reroute_edges(scene: FigureSceneGraph) -> None:
    nodes = scene.node_map()
    for edge in scene.edges:
        source = nodes.get(edge.source_id)
        target = nodes.get(edge.target_id)
        if not source or not target:
            continue
        direction = "down" if abs(source.bbox.center_x - target.bbox.center_x) < abs(source.bbox.center_y - target.bbox.center_y) else "right"
        edge.points = orthogonal_route(source.bbox, target.bbox, direction=direction)
        edge.source_port = "right" if direction == "right" else "bottom"
        edge.target_port = "left" if direction == "right" else "top"


def render_scene_to_png(scene: FigureSceneGraph, output_path: str | Path, scale: float = 1.0) -> Path:
    reroute_edges(scene)
    width = int(scene.width * scale)
    height = int(scene.height * scale)
    background = scene.style_tokens.palette["background"]
    image = Image.new("RGBA", (width, height), background)
    draw = ImageDraw.Draw(image)
    font = _load_font(max(14, int(scene.style_tokens.font_size * scale)))
    title_font = _load_font(max(16, int(scene.style_tokens.title_size * scale)))

    for group in sorted(scene.groups, key=lambda item: item.bbox.y):
        style = _group_style(scene, group)
        box = [
            group.bbox.x * scale,
            group.bbox.y * scale,
            group.bbox.right * scale,
            group.bbox.bottom * scale,
        ]
        draw.rounded_rectangle(
            box,
            radius=int(scene.style_tokens.corner_radius * scale),
            fill=style["fill"],
            outline=style["stroke"],
            width=max(1, int(scene.style_tokens.stroke_width * scale)),
        )
        draw.text(
            (group.bbox.x * scale + 18, group.bbox.y * scale + 12),
            group.label,
            fill=style["text"],
            font=title_font,
        )

    for edge in scene.edges:
        if len(edge.points) < 2:
            continue
        points = [(x * scale, y * scale) for x, y in edge.points]
        draw.line(
            points,
            fill=edge.style.get("stroke", scene.style_tokens.palette["accent"]),
            width=max(1, int(scene.style_tokens.arrow_width * scale)),
            joint="curve",
        )
        end = points[-1]
        prev = points[-2]
        angle = math.atan2(end[1] - prev[1], end[0] - prev[0])
        arrow_len = 14 * scale
        arrow_spread = math.pi / 7
        left = (
            end[0] - arrow_len * math.cos(angle - arrow_spread),
            end[1] - arrow_len * math.sin(angle - arrow_spread),
        )
        right = (
            end[0] - arrow_len * math.cos(angle + arrow_spread),
            end[1] - arrow_len * math.sin(angle + arrow_spread),
        )
        draw.polygon([end, left, right], fill=edge.style.get("stroke", scene.style_tokens.palette["accent"]))

    for node in sorted(scene.nodes, key=lambda item: item.z_index):
        style = _node_style(scene, node)
        box = [
            node.bbox.x * scale,
            node.bbox.y * scale,
            node.bbox.right * scale,
            node.bbox.bottom * scale,
        ]
        if node.kind == "background":
            draw.rectangle(box, fill=style["fill"])
            continue
        if node.kind == "icon":
            draw.ellipse(
                box,
                fill=style["fill"],
                outline=style["stroke"],
                width=max(1, int(scene.style_tokens.stroke_width * scale)),
            )
        else:
            draw.rounded_rectangle(
                box,
                radius=int(scene.style_tokens.corner_radius * scale),
                fill=style["fill"],
                outline=style["stroke"],
                width=max(1, int(scene.style_tokens.stroke_width * scale)),
            )
        if node.label:
            max_text_width = max(30, int(node.bbox.width * scale) - 24)
            lines = _wrap_text(draw, node.label, font, max_text_width)
            total_height = len(lines) * (font.size + 4)
            text_y = node.bbox.y * scale + max(12, (node.bbox.height * scale - total_height) / 2)
            for line in lines:
                bbox = draw.textbbox((0, 0), line, font=font)
                text_width = bbox[2] - bbox[0]
                text_x = node.bbox.x * scale + max(12, (node.bbox.width * scale - text_width) / 2)
                draw.text((text_x, text_y), line, fill=style["text"], font=font)
                text_y += font.size + 4

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    image.save(output_path)
    return output_path


def export_scene_to_svg(scene: FigureSceneGraph, output_path: str | Path) -> Path:
    reroute_edges(scene)
    parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{scene.width}" height="{scene.height}" viewBox="0 0 {scene.width} {scene.height}">',
        '<defs><marker id="arrow" markerWidth="12" markerHeight="12" refX="10" refY="6" orient="auto"><path d="M0,0 L12,6 L0,12 z" fill="#2563EB" /></marker></defs>',
        f'<rect width="{scene.width}" height="{scene.height}" fill="{scene.style_tokens.palette["background"]}" />',
    ]
    for group in scene.groups:
        style = _group_style(scene, group)
        parts.append(
            f'<rect x="{group.bbox.x}" y="{group.bbox.y}" width="{group.bbox.width}" height="{group.bbox.height}" '
            f'rx="{scene.style_tokens.corner_radius}" ry="{scene.style_tokens.corner_radius}" '
            f'fill="{style["fill"]}" stroke="{style["stroke"]}" stroke-width="{scene.style_tokens.stroke_width}" />'
        )
        parts.append(
            f'<text x="{group.bbox.x + 18}" y="{group.bbox.y + 34}" fill="{style["text"]}" '
            f'font-size="{scene.style_tokens.title_size}" font-family="{escape(scene.style_tokens.font_family)}">{escape(group.label)}</text>'
        )
    for edge in scene.edges:
        points = edge.points or []
        if len(points) < 2:
            continue
        d = " ".join(
            [f"M {points[0][0]} {points[0][1]}"] + [f"L {x} {y}" for x, y in points[1:]]
        )
        stroke = edge.style.get("stroke", scene.style_tokens.palette["accent"])
        parts.append(
            f'<path d="{d}" fill="none" stroke="{stroke}" stroke-width="{scene.style_tokens.arrow_width}" marker-end="url(#arrow)" />'
        )
    for node in scene.nodes:
        style = _node_style(scene, node)
        if node.kind == "background":
            continue
        if node.kind == "icon":
            parts.append(
                f'<ellipse cx="{node.bbox.center_x}" cy="{node.bbox.center_y}" rx="{node.bbox.width / 2}" ry="{node.bbox.height / 2}" '
                f'fill="{style["fill"]}" stroke="{style["stroke"]}" stroke-width="{scene.style_tokens.stroke_width}" />'
            )
        else:
            parts.append(
                f'<rect x="{node.bbox.x}" y="{node.bbox.y}" width="{node.bbox.width}" height="{node.bbox.height}" '
                f'rx="{scene.style_tokens.corner_radius}" ry="{scene.style_tokens.corner_radius}" '
                f'fill="{style["fill"]}" stroke="{style["stroke"]}" stroke-width="{scene.style_tokens.stroke_width}" />'
            )
        if node.label:
            parts.append(
                f'<text x="{node.bbox.center_x}" y="{node.bbox.center_y}" text-anchor="middle" dominant-baseline="middle" '
                f'fill="{style["text"]}" font-size="{scene.style_tokens.font_size}" font-family="{escape(scene.style_tokens.font_family)}">{escape(node.label)}</text>'
            )
    parts.append("</svg>")
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text("\n".join(parts), encoding="utf-8")
    return output_path


def _apply_shape_text(shape, text: str, font_size: int, color: RGBColor) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Emu(font_size * 12700)
    run.font.color.rgb = color


def _bbox_to_emu(box: BBox) -> tuple[int, int, int, int]:
    return (
        Emu(int(box.x * EMU_PER_PIXEL)),
        Emu(int(box.y * EMU_PER_PIXEL)),
        Emu(int(box.width * EMU_PER_PIXEL)),
        Emu(int(box.height * EMU_PER_PIXEL)),
    )


def export_scene_to_pptx(scene: FigureSceneGraph, output_path: str | Path) -> Path:
    reroute_edges(scene)
    presentation = Presentation()
    presentation.slide_width = Emu(scene.width * EMU_PER_PIXEL)
    presentation.slide_height = Emu(scene.height * EMU_PER_PIXEL)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    bg_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Emu(0),
        Emu(0),
        Emu(scene.width * EMU_PER_PIXEL),
        Emu(scene.height * EMU_PER_PIXEL),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _hex_to_rgb_color(scene.style_tokens.palette["background"])
    bg_shape.line.fill.background()

    for group in scene.groups:
        style = _group_style(scene, group)
        x, y, w, h = _bbox_to_emu(group.bbox)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb_color(style["fill"])
        shape.fill.transparency = 0.18
        shape.line.color.rgb = _hex_to_rgb_color(style["stroke"])
        shape.line.width = Emu(scene.style_tokens.stroke_width * 12700)
        _apply_shape_text(shape, group.label, scene.style_tokens.title_size, _hex_to_rgb_color(style["text"]))

    for node in scene.nodes:
        if node.kind == "background":
            continue
        style = _node_style(scene, node)
        x, y, w, h = _bbox_to_emu(node.bbox)
        shape_type = MSO_AUTO_SHAPE_TYPE.OVAL if node.kind == "icon" else MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        shape = slide.shapes.add_shape(shape_type, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb_color(style["fill"])
        shape.line.color.rgb = _hex_to_rgb_color(style["stroke"])
        shape.line.width = Emu(scene.style_tokens.stroke_width * 12700)
        if node.label:
            _apply_shape_text(shape, node.label, scene.style_tokens.font_size, _hex_to_rgb_color(style["text"]))

    for edge in scene.edges:
        points = edge.points or []
        if len(points) < 2:
            continue
        for idx in range(1, len(points)):
            x1, y1 = points[idx - 1]
            x2, y2 = points[idx]
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Emu(int(x1 * EMU_PER_PIXEL)),
                Emu(int(y1 * EMU_PER_PIXEL)),
                Emu(int(x2 * EMU_PER_PIXEL)),
                Emu(int(y2 * EMU_PER_PIXEL)),
            )
            connector.line.color.rgb = _hex_to_rgb_color(edge.style.get("stroke", scene.style_tokens.palette["accent"]))
            connector.line.width = Emu(scene.style_tokens.arrow_width * 12700)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return output_path
